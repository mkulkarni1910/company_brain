"""Generate the two SubStrateOS architecture decks (business + engineering).

Run: python3 docs/decks/build_decks.py
Output: docs/decks/SubStrateOS_Business.pptx, docs/decks/SubStrateOS_Engineering.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ---------------------------------------------------------------- theme
# Microsoft / Azure inspired palette.
INK = RGBColor(0x10, 0x1A, 0x2E)        # near-black navy text
NAVY = RGBColor(0x0B, 0x21, 0x46)       # deep navy
AZURE = RGBColor(0x00, 0x78, 0xD4)      # Azure blue
AZURE_LT = RGBColor(0x50, 0xA0, 0xE6)   # light azure
CYAN = RGBColor(0x00, 0xB7, 0xC3)       # teal accent
TEAL_DK = RGBColor(0x00, 0x6A, 0x71)
SLATE = RGBColor(0x5B, 0x6B, 0x82)      # secondary text
MIST = RGBColor(0xEC, 0xF2, 0xFB)       # pale panel fill
MIST2 = RGBColor(0xDD, 0xE9, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0x88, 0x4A)      # "shipped" green
AMBER = RGBColor(0xB7, 0x6E, 0x00)      # roadmap amber
LINE = RGBColor(0xC8, 0xD4, 0xE3)

FONT = "Segoe UI"
FONT_LT = "Segoe UI Light"

# 16:9
EMU = 914400
SW = int(13.333 * EMU)
SH = int(7.5 * EMU)


def _in(v: float) -> int:
    return int(v * EMU)


def new_deck() -> Presentation:
    p = Presentation()
    p.slide_width = SW
    p.slide_height = SH
    return p


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _line(shape, color, w=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(w)


def _no_shadow(shape):
    el = shape._element.spPr
    # remove inherited preset shadow
    existing = el.find(qn("a:effectLst"))
    if existing is None:
        el.append(el.makeelement(qn("a:effectLst"), {}))


def rect(slide, x, y, w, h, fill=None, line_color=None, line_w=1.0,
         shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, _in(x), _in(y), _in(w), _in(h))
    if fill is None:
        sp.fill.background()
    else:
        _set_fill(sp, fill)
    if line_color is None:
        _no_line(sp)
    else:
        _line(sp, line_color, line_w)
    _no_shadow(sp)
    sp.text_frame.word_wrap = True
    return sp


def txt(slide, x, y, w, h, text, size=18, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0,
        space_after=2):
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        para.space_after = Pt(space_after)
        para.space_before = Pt(0)
        r = para.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(size)
        f.name = font
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def fill_text(shape, text, size=14, color=INK, bold=False, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE, font=FONT, line_spacing=1.0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        r = para.add_run()
        r.text = ln
        f = r.font
        f.size = Pt(size)
        f.name = font
        f.bold = bold
        f.color.rgb = color
    return shape


def bg(slide, color=WHITE):
    rect(slide, -0.05, -0.05, 13.43, 7.6, fill=color)


def header(slide, title, kicker=None, accent=AZURE):
    """Standard content-slide header: accent bar + title + kicker."""
    rect(slide, 0.0, 0.0, 0.22, 7.5, fill=accent)
    if kicker:
        txt(slide, 0.7, 0.42, 11.5, 0.4, kicker.upper(), size=12.5, color=accent,
            bold=True)
        txt(slide, 0.7, 0.74, 11.9, 0.9, title, size=29, color=NAVY, bold=True,
            font=FONT_LT)
    else:
        txt(slide, 0.7, 0.55, 11.9, 0.9, title, size=30, color=NAVY, bold=True,
            font=FONT_LT)
    rect(slide, 0.72, 1.62, 1.1, 0.045, fill=CYAN)


def footer(slide, idx, label):
    txt(slide, 0.7, 7.05, 8.0, 0.3, label, size=9, color=SLATE)
    txt(slide, 11.8, 7.05, 0.9, 0.3, f"{idx:02d}", size=9, color=SLATE,
        align=PP_ALIGN.RIGHT)


def chip(slide, x, y, w, h, label, fill=MIST, text_color=NAVY, line_color=LINE,
         size=12.5, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = rect(slide, x, y, w, h, fill=fill, line_color=line_color, line_w=1.0,
              shape=shape)
    try:
        sp.adjustments[0] = 0.12
    except Exception:
        pass
    fill_text(sp, label, size=size, color=text_color, bold=bold)
    return sp


def arrow(slide, x, y, w, h=0.0, color=AZURE, w_pt=2.5, vertical=False):
    """Connector line with arrowhead."""
    from pptx.enum.shapes import MSO_CONNECTOR
    if vertical:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _in(x), _in(y),
                                          _in(x), _in(y + h))
    else:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, _in(x), _in(y),
                                          _in(x + w), _in(y))
    conn.line.color.rgb = color
    conn.line.width = Pt(w_pt)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"),
                          {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


# ================================================================ shared slides
def title_slide(prs, deck_kind, subtitle, audience):
    s = blank(prs)
    bg(s, NAVY)
    # diagonal accent band
    band = rect(s, -1, 4.7, 16, 3.2, fill=RGBColor(0x09, 0x1A, 0x38))
    band.rotation = -4
    rect(s, 0, 0, 13.333, 0.16, fill=AZURE)
    # eyebrow
    txt(s, 0.9, 1.25, 11, 0.5, "MICROSOFT  ·  AZURE AI", size=14, color=AZURE_LT,
        bold=True)
    txt(s, 0.85, 1.95, 11.6, 1.6, "SubStrateOS", size=66, color=WHITE, bold=True,
        font=FONT_LT)
    txt(s, 0.9, 3.25, 11.5, 0.8,
        "Enterprise intelligence layer on Microsoft Azure", size=24,
        color=AZURE_LT, font=FONT_LT)
    # rule
    rect(s, 0.92, 4.15, 2.2, 0.05, fill=CYAN)
    txt(s, 0.9, 4.45, 11.5, 0.8, subtitle, size=20, color=WHITE)
    txt(s, 0.9, 6.55, 11.5, 0.5, audience, size=13, color=AZURE_LT, bold=True)
    return s


def logo_wall(slide, x, y, items, cols=4, cw=2.7, ch=0.95, gap=0.18,
              fill=WHITE, line_color=LINE, text_color=NAVY, accent=AZURE):
    """Grid of Azure service 'logo' tiles (two-line: bold name + sub)."""
    for i, (name, sub) in enumerate(items):
        r = i // cols
        c = i % cols
        cx = x + c * (cw + gap)
        cy = y + r * (ch + gap)
        tile = rect(slide, cx, cy, cw, ch, fill=fill, line_color=line_color,
                    line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        try:
            tile.adjustments[0] = 0.09
        except Exception:
            pass
        # accent square
        rect(slide, cx + 0.16, cy + ch / 2 - 0.18, 0.36, 0.36, fill=accent,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb = slide.shapes.add_textbox(_in(cx + 0.68), _in(cy + 0.12),
                                      _in(cw - 0.78), _in(ch - 0.24))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = name
        r0.font.size = Pt(13)
        r0.font.bold = True
        r0.font.name = FONT
        r0.font.color.rgb = text_color
        p1 = tf.add_paragraph()
        r1 = p1.add_run()
        r1.text = sub
        r1.font.size = Pt(9.5)
        r1.font.name = FONT
        r1.font.color.rgb = SLATE


def bullets(slide, x, y, w, items, size=16, gap=0.62, color=INK, marker_color=AZURE,
            bold_lead=True):
    for i, item in enumerate(items):
        cy = y + i * gap
        # marker
        rect(slide, x, cy + 0.09, 0.16, 0.16, fill=marker_color,
             shape=MSO_SHAPE.OVAL)
        if isinstance(item, tuple):
            lead, rest = item
            tb = slide.shapes.add_textbox(_in(x + 0.34), _in(cy - 0.02),
                                          _in(w - 0.34), _in(gap))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = lead + "  "
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.name = FONT
            r1.font.color.rgb = NAVY
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.name = FONT
            r2.font.color.rgb = color
        else:
            txt(slide, x + 0.34, cy - 0.02, w - 0.34, gap, item, size=size,
                color=color)


def section_caption(slide, x, y, w, text, color=SLATE, size=13):
    txt(slide, x, y, w, 0.5, text, size=size, color=color, italic=True)


# ================================================================ BUSINESS DECK
def build_business():
    prs = new_deck()
    LABEL = "SubStrateOS · Microsoft Hackathon Review"

    # 1 — Title
    title_slide(prs, "business",
                "Find any answer across your company — grounded, personalized, secure.",
                "FOR THE HACKATHON REVIEW COMMITTEE")

    # 2 — Problem
    s = blank(prs); bg(s)
    header(s, "Company knowledge is everywhere — and nowhere",
           kicker="The problem")
    txt(s, 0.72, 1.95, 7.0, 1.4,
        "An employee's answer lives scattered across SharePoint files, Outlook "
        "mail, Teams chats, and a dozen apps. Finding it means asking around, "
        "guessing keywords, and re-doing work that already exists.",
        size=18, color=SLATE, line_spacing=1.15)
    cards = [
        ("Wasted time", "Knowledge workers spend ~20% of the week just searching for information."),
        ("Wrong answers", "Stale docs and dead links lead to decisions on outdated facts."),
        ("Locked-in silos", "Each tool searches only itself — no single, permission-aware view."),
    ]
    cx = 0.72
    for name, body in cards:
        c = rect(s, cx, 3.7, 3.78, 2.55, fill=MIST, line_color=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, cx, 3.7, 3.78, 0.12, fill=AZURE)
        txt(s, cx + 0.3, 4.05, 3.2, 0.5, name, size=18, color=NAVY, bold=True)
        txt(s, cx + 0.3, 4.65, 3.2, 1.4, body, size=13.5, color=SLATE,
            line_spacing=1.12)
        cx += 4.0
    footer(s, 2, LABEL)

    # 3 — Solution
    s = blank(prs); bg(s)
    header(s, "One question box over everything you're allowed to see",
           kicker="The solution")
    txt(s, 0.72, 1.95, 11.8, 1.0,
        "SubStrateOS reads across your Microsoft 365 world and answers in plain "
        "language — every answer backed by citations to the real source.",
        size=18, color=SLATE, line_spacing=1.15)
    # mini flow
    chip(s, 0.9, 3.4, 2.7, 1.1, "Ask a question\nin plain English", fill=AZURE,
         text_color=WHITE, line_color=AZURE, size=14)
    chip(s, 5.0, 3.4, 3.0, 1.1, "SubStrateOS\nsearches what YOU can access",
         fill=NAVY, text_color=WHITE, line_color=NAVY, size=13.5)
    chip(s, 9.4, 3.4, 3.0, 1.1, "Grounded answer\n+ citations", fill=CYAN,
         text_color=WHITE, line_color=CYAN, size=14)
    arrow(s, 3.7, 3.95, 1.2)
    arrow(s, 8.1, 3.95, 1.2)
    txt(s, 0.9, 5.1, 11.6, 1.0,
        "“What is our PTO policy?”  →  a clear answer with a link to the exact HR "
        "doc it came from. No hallucinations: if the answer isn't in your "
        "content, it says so.", size=15, color=NAVY, italic=True, line_spacing=1.15)
    footer(s, 3, LABEL)

    # 4 — What it does (3 capabilities)
    s = blank(prs); bg(s)
    header(s, "Three things that make it trustworthy", kicker="What it does")
    caps = [
        ("Grounded", "Every answer cites its sources.",
         "Answers come only from your real documents, with links to verify — never invented."),
        ("Personalized", "Tuned to who you are.",
         "Ranks results using your team, your permissions, and what your org actually engages with."),
        ("Always fresh", "Live, not last-week.",
         "For 'today / right now' questions it reaches into Microsoft 365 live, at the moment you ask."),
    ]
    cx = 0.72
    for name, lead, body in caps:
        c = rect(s, cx, 2.2, 3.85, 3.9, fill=WHITE, line_color=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, cx + 0.32, 2.55, 0.7, 0.7, fill=AZURE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, cx + 0.32, 3.5, 3.2, 0.5, name, size=21, color=NAVY, bold=True)
        txt(s, cx + 0.32, 4.05, 3.2, 0.5, lead, size=14, color=AZURE, bold=True)
        txt(s, cx + 0.32, 4.6, 3.25, 1.4, body, size=13.5, color=SLATE,
            line_spacing=1.15)
        cx += 4.04
    footer(s, 4, LABEL)

    # 5 — Simple architecture (non-techie)
    s = blank(prs); bg(s)
    header(s, "How it works — the simple view", kicker="Architecture at a glance")
    # user
    chip(s, 1.0, 2.7, 2.4, 1.3, "You\nask a question", fill=AZURE, text_color=WHITE,
         line_color=AZURE, size=15)
    # brain
    brain = rect(s, 4.5, 2.4, 4.3, 1.9, fill=NAVY,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    fill_text(brain, "SubStrateOS\nintelligence layer", size=18, color=WHITE,
              bold=True)
    txt(s, 4.5, 4.0, 4.3, 0.4, "(understands · finds · ranks · explains)",
        size=11.5, color=AZURE_LT, align=PP_ALIGN.CENTER)
    # data sources
    srcs = ["SharePoint\n& files", "Outlook\nmail & calendar", "Teams\nconversations"]
    sy = 1.95
    for i, name in enumerate(srcs):
        c = chip(s, 10.0, sy, 2.7, 1.1, name, fill=MIST, text_color=NAVY,
                 line_color=LINE, size=13)
        arrow(s, 8.8, sy + 0.55, 1.2, color=CYAN)
        sy += 1.45
    arrow(s, 3.4, 3.05, 1.1)
    # return arrow
    arrow(s, 4.5, 5.0, -3.5, color=GREEN)  # not visible direction issue; draw label
    txt(s, 1.0, 5.05, 3.4, 0.5, "Answer + citations", size=12.5, color=GREEN,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, 0.72, 6.05, 12, 0.6,
        "Everything runs on Microsoft Azure. You only ever see content you already "
        "have permission to open.", size=14, color=SLATE, italic=True)
    footer(s, 5, LABEL)

    # 6 — Built on Microsoft (logo wall)
    s = blank(prs); bg(s)
    header(s, "Built on the Microsoft Azure stack", kicker="Why it's enterprise-ready")
    txt(s, 0.72, 1.82, 11.8, 0.6,
        "Every layer your IT team governs — search, data, identity, and hosting — "
        "is a first-party Azure service. No exotic dependencies.", size=15,
        color=SLATE)
    items = [
        ("Azure AI Search", "Hybrid retrieval"),
        ("Azure OpenAI", "Embeddings + vectors"),
        ("Cosmos DB", "People graph"),
        ("Azure Data Explorer", "Activity signals"),
        ("Microsoft Graph", "Live M365 data"),
        ("Microsoft Entra ID", "Sign-in & identity"),
        ("Azure Container Apps", "Hosting"),
        ("Azure Monitor", "Observability"),
    ]
    logo_wall(s, 0.9, 2.6, items, cols=4, cw=2.85, ch=1.2, gap=0.2)
    note = rect(s, 0.9, 5.7, 11.65, 0.75, fill=MIST, line_color=LINE,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 1.2, 5.83, 11.1, 0.55,
        "Answer generation runs on Google Gemini 2.5 Pro through a pluggable LLM "
        "layer — swappable for Azure OpenAI with no architecture change.",
        size=12.5, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 6, LABEL)

    # 7 — Security & trust
    s = blank(prs); bg(s)
    header(s, "You only ever see what you're permitted to", kicker="Security & trust",
           accent=TEAL_DK)
    txt(s, 0.72, 1.95, 7.2, 1.6,
        "Permissions aren't an afterthought — they're enforced twice, "
        "independently. If the permission service is ever unreachable, the system "
        "shows nothing rather than risk a leak.",
        size=18, color=SLATE, line_spacing=1.2)
    pts = [
        ("Permission-stamped", "every document carries its access list from day one"),
        ("Re-checked live", "access is verified again the moment you ask"),
        ("Fail-closed", "no permission confirmation → no result, ever"),
        ("Microsoft Entra ID", "enterprise single sign-on, no separate passwords"),
    ]
    bullets(s, 0.72, 3.9, 7.2, pts, size=15, gap=0.62, marker_color=TEAL_DK)
    # shield panel
    panel = rect(s, 8.6, 2.2, 4.0, 4.0, fill=MIST, line_color=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 10.15, 2.6, 0.9, 0.9, fill=TEAL_DK, shape=MSO_SHAPE.HEXAGON)
    txt(s, 8.7, 3.7, 3.8, 0.5, "Double enforcement", size=17, color=NAVY,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, 8.8, 4.25, 3.6, 1.8,
        "Index-time access stamp\n+\nQuery-time access re-check\n=\nzero "
        "over-sharing by design", size=13.5, color=SLATE, align=PP_ALIGN.CENTER,
        line_spacing=1.3)
    footer(s, 7, LABEL)

    # 8 — Personalization differentiator
    s = blank(prs); bg(s)
    header(s, "Same question, the right answer for each person",
           kicker="Our differentiator")
    txt(s, 0.72, 1.9, 11.8, 0.9,
        "A generic search returns one ranking for everyone. SubStrateOS reorders "
        "results around the person asking — their team, their work, their access.",
        size=16, color=SLATE, line_spacing=1.15)
    # two personas
    def persona(x, name, role, color, ranking):
        rect(s, x, 3.1, 5.6, 3.1, fill=WHITE, line_color=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, 3.1, 5.6, 0.6, fill=color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + 0.3, 3.16, 5.0, 0.5, f"{name} — {role}", size=14, color=WHITE,
            bold=True)
        for i, r in enumerate(ranking):
            yy = 3.95 + i * 0.62
            rect(s, x + 0.35, yy, 0.42, 0.42, fill=color,
                 shape=MSO_SHAPE.OVAL)
            txt(s, x + 0.36, yy + 0.04, 0.4, 0.4, str(i + 1), size=13,
                color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            txt(s, x + 0.95, yy + 0.02, 4.5, 0.5, r, size=13.5, color=INK)
    persona(0.72, "Maya", "Engineering", AZURE,
            ["On-call runbook", "Service alert guide", "Q3 ARR target"])
    persona(7.0, "Raj", "Finance", CYAN,
            ["Q3 ARR target", "Revenue model", "On-call runbook"])
    txt(s, 0.72, 6.35, 12, 0.5,
        "Identical query — ranking shaped by role, relationships, and real "
        "engagement.", size=13, color=SLATE, italic=True)
    footer(s, 8, LABEL)

    # 9 — What's shipped
    s = blank(prs); bg(s)
    header(s, "Already built and running on Azure", kicker="Proof of execution",
           accent=GREEN)
    rows = [
        ("Phase 1", "Grounded Q&A with citations", "Live", GREEN),
        ("Phase 2a", "People graph + personalized ranking", "Live", GREEN),
        ("Phase 2b", "Activity signals from real engagement", "Live", GREEN),
        ("Phase 3", "Live Fetch from Microsoft Graph", "Live", GREEN),
        ("Phase 4", "Smarter planning, ranking & ACL gating", "Live", GREEN),
        ("Connectors", "SharePoint · Outlook · Teams", "Code-complete", AMBER),
    ]
    y = 2.3
    for ph, desc, status, col in rows:
        rect(s, 0.9, y, 11.6, 0.62, fill=MIST if (rows.index((ph,desc,status,col))%2==0) else WHITE,
             line_color=LINE)
        txt(s, 1.15, y + 0.13, 2.2, 0.4, ph, size=15, color=NAVY, bold=True)
        txt(s, 3.5, y + 0.14, 6.8, 0.4, desc, size=14, color=INK)
        badge = rect(s, 10.4, y + 0.11, 1.85, 0.4, fill=col,
                     shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        fill_text(badge, status, size=11.5, color=WHITE, bold=True)
        y += 0.66
    txt(s, 0.9, 6.35, 12, 0.5,
        "Quality is measured, not claimed: golden-question eval at 1.0 recall / "
        "1.0 MRR on the Phase-1 baseline.", size=13, color=SLATE, italic=True)
    footer(s, 9, LABEL)

    # 10 — Roadmap
    s = blank(prs); bg(s)
    header(s, "Where it goes next", kicker="Roadmap", accent=AMBER)
    steps = [
        ("Now", "Grounded Q&A, personalization, live fetch — on Azure", GREEN),
        ("Next", "Full connector rollout + Model Context Protocol for AI agents", AZURE),
        ("Then", "Per-user delegated access, API gateway, multi-tenant isolation", AMBER),
    ]
    x = 0.95
    for i, (when, body, col) in enumerate(steps):
        rect(s, x, 2.9, 3.7, 2.6, fill=WHITE, line_color=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, 2.9, 3.7, 0.7, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + 0.3, 2.98, 3.2, 0.5, when, size=18, color=WHITE, bold=True)
        txt(s, x + 0.3, 3.85, 3.15, 1.5, body, size=14, color=SLATE,
            line_spacing=1.2)
        if i < 2:
            arrow(s, x + 3.7, 4.2, 0.42, color=SLATE)
        x += 4.12
    txt(s, 0.95, 5.9, 11.5, 0.6,
        "Built monolith-first with clean module boundaries — each component can "
        "graduate into its own service with no rewrite.", size=14, color=SLATE,
        italic=True)
    footer(s, 10, LABEL)

    # 11 — Closing / impact
    s = blank(prs); bg(s, NAVY)
    rect(s, 0, 0, 13.333, 0.16, fill=AZURE)
    txt(s, 0.9, 1.4, 11.5, 0.5, "THE IMPACT", size=14, color=AZURE_LT, bold=True)
    txt(s, 0.85, 2.0, 11.6, 1.6,
        "Every employee's questions,\nanswered from your own knowledge —\nsecurely.",
        size=38, color=WHITE, bold=True, font=FONT_LT, line_spacing=1.05)
    metrics = [
        ("5", "Azure AI services"),
        ("4", "phases live"),
        ("1.0", "eval recall & MRR"),
        ("0", "over-sharing by design"),
    ]
    x = 0.95
    for val, lab in metrics:
        rect(s, x, 4.5, 2.8, 1.5, fill=RGBColor(0x0E, 0x29, 0x4F),
             line_color=AZURE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x, 4.7, 2.8, 0.7, val, size=34, color=CYAN, bold=True,
            align=PP_ALIGN.CENTER, font=FONT_LT)
        txt(s, x + 0.15, 5.5, 2.5, 0.5, lab, size=12, color=WHITE,
            align=PP_ALIGN.CENTER)
        x += 3.0
    txt(s, 0.9, 6.6, 11.5, 0.5, "SubStrateOS  ·  built on Microsoft Azure",
        size=13, color=AZURE_LT, bold=True)

    out = os.path.join(os.path.dirname(__file__), "SubStrateOS_Business.pptx")
    prs.save(out)
    return out


# ================================================================ ENGINEERING DECK
def build_engineering():
    prs = new_deck()
    LABEL = "SubStrateOS · Architecture deep dive"

    # 1 — Title
    title_slide(prs, "eng", "Architecture deep dive — a Glean-class Zone 4 "
                "intelligence layer on Azure.",
                "FOR THE TECHNICAL REVIEW PANEL")

    # 2 — System topology
    s = blank(prs); bg(s)
    header(s, "System topology", kicker="End to end")
    # client
    chip(s, 0.8, 2.15, 2.7, 1.0, "Next.js 14 Web Chat\nMSAL · Entra SSO",
         fill=AZURE, text_color=WHITE, line_color=AZURE, size=13)
    # ingress
    chip(s, 0.8, 3.7, 2.7, 0.9, "Azure Container Apps\ningress (HTTPS + JWT)",
         fill=NAVY, text_color=WHITE, line_color=NAVY, size=12.5)
    arrow(s, 2.15, 3.15, 0.0, h=0.55, vertical=True)
    # API box
    api = rect(s, 4.4, 1.9, 5.0, 4.6, fill=MIST, line_color=AZURE, line_w=1.5,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 4.6, 2.0, 4.6, 0.4, "substrateos-api  ·  FastAPI (Python 3.12)",
        size=13, color=NAVY, bold=True)
    txt(s, 4.6, 2.4, 4.6, 0.35, "Zone 4 intelligence layer — modular monolith",
        size=10.5, color=SLATE, italic=True)
    mods = [
        "Semantic Kernel orchestrator  /query",
        "Hybrid retriever  (RRF)",
        "Personalized ranker  (multi-signal)",
        "Live Fetch  (Graph /search)",
        "ACL recheck  (fail-closed)",
        "Generation  (Gemini 2.5 Pro + citations)",
        "Admin · ingest · connectors · MCP",
    ]
    my = 2.85
    for m in mods:
        b = rect(s, 4.65, my, 4.5, 0.46, fill=WHITE, line_color=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, 4.85, my + 0.08, 4.2, 0.35, m, size=11.5, color=INK)
        my += 0.5
    arrow(s, 3.5, 4.15, 0.9)
    # backing services column
    svcs = [("Azure AI Search", CYAN), ("Cosmos DB (Gremlin)", CYAN),
            ("Azure Data Explorer", CYAN), ("Azure OpenAI · embeddings", CYAN),
            ("Gemini 2.5 Pro / Flash", AMBER), ("Redis · Key Vault", CYAN),
            ("Microsoft Graph", CYAN)]
    sy = 1.9
    for name, dot in svcs:
        c = rect(s, 10.3, sy, 2.55, 0.56, fill=WHITE, line_color=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, 10.42, sy + 0.14, 0.26, 0.26, fill=dot,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, 10.78, sy + 0.12, 2.05, 0.4, name, size=11, color=NAVY, bold=True)
        sy += 0.7
    arrow(s, 9.4, 4.2, 0.85, color=CYAN)
    txt(s, 0.72, 6.7, 12, 0.4,
        "DefaultAzureCredential → user-assigned managed identity (no secrets in "
        "code). OpenTelemetry → Application Insights, one trace per request.",
        size=11.5, color=SLATE, italic=True)
    footer(s, 2, LABEL)

    # 3 — Request lifecycle
    s = blank(prs); bg(s)
    header(s, "Request lifecycle — POST /query", kicker="Orchestration flow")
    steps = [
        ("1", "Resolve user", "Entra JWT / Easy-Auth / PAT → tenant + groups"),
        ("2", "Plan", "Semantic Kernel classifies intent (Gemini 2.5 Flash)"),
        ("3", "Retrieve", "Hybrid search + People proximity + Activity scores"),
        ("4", "Live Fetch?", "Freshness trigger → Microsoft Graph /search merge"),
        ("5", "ACL recheck", "Re-verify access per candidate · fail-closed"),
        ("6", "Rank", "RRF + weighted signals → ordered context"),
        ("7", "Generate", "Gemini 2.5 Pro grounded answer with inline citations"),
    ]
    y = 2.1
    for num, title, desc in steps:
        rect(s, 0.9, y, 0.62, 0.62, fill=AZURE, shape=MSO_SHAPE.OVAL)
        txt(s, 0.9, y + 0.11, 0.62, 0.4, num, size=20, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER)
        b = rect(s, 1.75, y, 10.7, 0.62, fill=MIST, line_color=LINE,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, 2.0, y + 0.14, 2.6, 0.4, title, size=14.5, color=NAVY, bold=True)
        txt(s, 4.7, y + 0.15, 7.5, 0.4, desc, size=13, color=SLATE)
        if num != "7":
            arrow(s, 1.21, y + 0.62, 0.0, h=0.085, vertical=True, color=AZURE_LT,
                  w_pt=2)
        y += 0.7
    footer(s, 3, LABEL)

    # 4 — Three pillars
    s = blank(prs); bg(s)
    header(s, "Three pillars feed every answer", kicker="Knowledge model")
    pillars = [
        ("Content", "Azure AI Search", AZURE,
         ["Hybrid: BM25 + vector", "text-embedding-3-large", "ACL-stamped chunks",
          "Semantic ranking"]),
        ("People", "Cosmos DB · Gremlin", CYAN,
         ["Org/collab graph", "Author ↔ asker proximity", "Group membership",
          "Drives personalization"]),
        ("Activity", "Azure Data Explorer", TEAL_DK,
         ["Click / dwell events", "Per-type engagement", "Recency decay",
          "Real signal, not vibes"]),
    ]
    x = 0.72
    for name, svc, col, feats in pillars:
        rect(s, x, 2.1, 3.85, 4.3, fill=WHITE, line_color=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, 2.1, 3.85, 0.95, fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + 0.3, 2.2, 3.3, 0.5, name, size=21, color=WHITE, bold=True)
        txt(s, x + 0.3, 2.68, 3.3, 0.35, svc, size=12, color=WHITE)
        for i, f in enumerate(feats):
            yy = 3.35 + i * 0.66
            rect(s, x + 0.35, yy + 0.07, 0.13, 0.13, fill=col,
                 shape=MSO_SHAPE.OVAL)
            txt(s, x + 0.62, yy, 3.0, 0.5, f, size=13, color=INK)
        x += 4.04
    txt(s, 0.72, 6.55, 12, 0.4,
        "Each pillar is an independent module behind a stable interface — "
        "swappable and individually testable.", size=12, color=SLATE, italic=True)
    footer(s, 4, LABEL)

    # 5 — Retrieval + ranker
    s = blank(prs); bg(s)
    header(s, "Hybrid retrieval + personalized ranker", kicker="Ranking")
    txt(s, 0.72, 1.85, 11.8, 0.5,
        "Candidates fused with Reciprocal Rank Fusion, then re-scored by a "
        "weighted multi-signal ranker — same query, different order per user.",
        size=14.5, color=SLATE, line_spacing=1.1)
    # signal boxes feeding ranker
    sigs = [("Content\nrelevance", AZURE), ("People\nproximity", CYAN),
            ("Activity\nengagement", TEAL_DK), ("Recency\nsignal", AMBER)]
    x = 0.9
    for name, col in sigs:
        c = rect(s, x, 2.75, 2.5, 1.15, fill=WHITE, line_color=col, line_w=1.5,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        fill_text(c, name, size=13.5, color=NAVY, bold=True)
        arrow(s, x + 1.25, 3.9, 0.0, h=0.55, vertical=True, color=col, w_pt=2)
        x += 3.0
    ranker = rect(s, 2.4, 4.65, 8.5, 0.95, fill=NAVY,
                  shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    fill_text(ranker, "Personalized Ranker   —   RRF ⊕ weighted signal blend",
              size=16, color=WHITE, bold=True)
    arrow(s, 6.65, 5.6, 0.0, h=0.4, vertical=True, color=GREEN, w_pt=2.5)
    out = rect(s, 4.0, 6.05, 5.3, 0.7, fill=GREEN,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    fill_text(out, "Ranked, grounded context → generation", size=14, color=WHITE,
              bold=True)
    footer(s, 5, LABEL)

    # 6 — Live Fetch
    s = blank(prs); bg(s)
    header(s, "Live Fetch — freshness at query time", kicker="Real-time path")
    txt(s, 0.72, 1.9, 11.8, 1.0,
        "Most queries answer from the index. 'Today / right now' queries trip a "
        "freshness trigger that reaches into Microsoft 365 live and merges the "
        "results into ranking — so answers are never stale.",
        size=16, color=SLATE, line_spacing=1.2)
    # decision flow
    chip(s, 0.9, 3.5, 2.6, 1.0, "Query", fill=AZURE, text_color=WHITE,
         line_color=AZURE, size=15)
    dia = rect(s, 4.2, 3.35, 2.6, 1.3, fill=AMBER, shape=MSO_SHAPE.DIAMOND)
    fill_text(dia, "Freshness\ntrigger?", size=13, color=WHITE, bold=True)
    arrow(s, 3.5, 4.0, 0.7)
    # no path
    chip(s, 7.6, 2.35, 4.0, 0.95, "Index only  →  Azure AI Search",
         fill=MIST, text_color=NAVY, line_color=LINE, size=13)
    arrow(s, 6.8, 3.7, 0.8, color=SLATE)
    txt(s, 6.95, 3.1, 0.7, 0.3, "no", size=11, color=SLATE, bold=True)
    # yes path
    chip(s, 7.6, 4.8, 4.0, 0.95, "Live  →  Microsoft Graph /search",
         fill=NAVY, text_color=WHITE, line_color=NAVY, size=13)
    arrow(s, 5.5, 4.65, 0.0, h=0.6, vertical=True, color=AMBER)
    arrow(s, 5.5, 5.25, 2.1, color=AMBER)
    txt(s, 5.6, 4.75, 0.8, 0.3, "yes", size=11, color=AMBER, bold=True)
    txt(s, 0.72, 6.2, 12, 0.6,
        "Today: DefaultAzureCredential (single identity), heuristic + LLM trigger. "
        "Phase 5: per-user on-behalf-of (OBO) delegation.", size=13, color=SLATE,
        italic=True)
    footer(s, 6, LABEL)

    # 7 — Double-enforcement ACL
    s = blank(prs); bg(s)
    header(s, "Double-enforcement ACLs — fail-closed", kicker="Security model",
           accent=TEAL_DK)
    # two-stage
    a = rect(s, 0.9, 2.4, 5.4, 1.7, fill=MIST, line_color=TEAL_DK, line_w=1.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 1.15, 2.55, 5.0, 0.4, "① Index-time", size=15, color=TEAL_DK, bold=True)
    txt(s, 1.15, 2.95, 5.0, 1.0,
        "Ingest resolves each document's ACL and stamps allowed principals onto "
        "every chunk in Azure AI Search.", size=13, color=SLATE, line_spacing=1.1)
    b = rect(s, 0.9, 4.35, 5.4, 1.7, fill=MIST, line_color=TEAL_DK, line_w=1.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 1.15, 4.5, 5.0, 0.4, "② Query-time", size=15, color=TEAL_DK, bold=True)
    txt(s, 1.15, 4.9, 5.0, 1.1,
        "Every candidate is re-checked against the live ACL store at answer time — "
        "filter applied before generation.", size=13, color=SLATE, line_spacing=1.1)
    arrow(s, 3.6, 4.1, 0.0, h=0.25, vertical=True, color=TEAL_DK)
    # fail-closed panel
    p = rect(s, 6.8, 2.4, 5.7, 3.65, fill=NAVY,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 8.95, 2.75, 1.0, 1.0, fill=TEAL_DK, shape=MSO_SHAPE.HEXAGON)
    txt(s, 6.95, 3.95, 5.4, 0.5, "Fail-closed by default", size=19, color=WHITE,
        bold=True, align=PP_ALIGN.CENTER)
    txt(s, 7.1, 4.55, 5.1, 1.4,
        "ACL store unreachable  →  candidate dropped, never shown.\n"
        "An ACL-freshness gate forces re-resolution when stamps are stale.\n"
        "Over-sharing is structurally impossible.",
        size=13, color=AZURE_LT, align=PP_ALIGN.CENTER, line_spacing=1.25)
    footer(s, 7, LABEL)

    # 8 — Generation + eval
    s = blank(prs); bg(s)
    header(s, "Grounded generation, measured quality", kicker="Answer + evidence")
    txt(s, 0.72, 1.95, 5.6, 0.5, "Generation", size=17, color=NAVY, bold=True)
    bullets(s, 0.72, 2.55, 5.7, [
        ("Gemini 2.5 Pro", "grounded on ranked context only"),
        ("Inline citations", "every claim links to its source chunk"),
        ("Refuses out-of-corpus", "no grounding → explicit 'I don't know'"),
        ("Gemini 2.5 Flash planner", "cheap intent classification step"),
    ], size=13.5, gap=0.66)
    txt(s, 6.7, 1.95, 5.6, 0.5, "Eval harness (day one)", size=17, color=NAVY,
        bold=True)
    bullets(s, 6.7, 2.55, 5.7, [
        ("Golden Q&A regression", "quality changes are measurable"),
        ("A/B config harness", "ranker tuning is data-driven"),
        ("Retrieval metrics", "recall@10, MRR@10 gates"),
    ], size=13.5, gap=0.66)
    # metric tiles
    tiles = [("1.0", "recall@10"), ("1.0", "MRR@10"), ("10", "golden Qs"),
             ("~6s", "warm latency")]
    x = 0.9
    for val, lab in tiles:
        rect(s, x, 4.9, 2.75, 1.3, fill=MIST, line_color=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x, 5.05, 2.75, 0.6, val, size=30, color=AZURE, bold=True,
            align=PP_ALIGN.CENTER, font=FONT_LT)
        txt(s, x, 5.7, 2.75, 0.4, lab, size=12, color=SLATE, align=PP_ALIGN.CENTER)
        x += 2.95
    txt(s, 0.72, 6.5, 12, 0.4,
        "Phase-1 baseline, 2026-05-29 — quality is regression-tested, not vibes.",
        size=12, color=SLATE, italic=True)
    footer(s, 8, LABEL)

    # 9 — Connectors & extensibility
    s = blank(prs); bg(s)
    header(s, "Connectors & extensibility", kicker="Reach & integration")
    # connectors row
    txt(s, 0.72, 1.95, 11, 0.4, "Source connectors (Microsoft Graph)", size=15,
        color=NAVY, bold=True)
    conns = [("SharePoint", "files + sites, ACL-aware"),
             ("Outlook", "mail + calendar, per-participant ACL"),
             ("Teams", "channels + private-channel ACL")]
    x = 0.9
    for name, sub in conns:
        rect(s, x, 2.45, 3.75, 1.3, fill=WHITE, line_color=AZURE, line_w=1.3,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + 0.3, 2.62, 3.2, 0.5, name, size=16, color=AZURE, bold=True)
        txt(s, x + 0.3, 3.1, 3.2, 0.6, sub, size=12, color=SLATE,
            line_spacing=1.05)
        x += 3.95
    # interfaces row
    txt(s, 0.72, 4.15, 11, 0.4, "Programmatic surfaces", size=15, color=NAVY,
        bold=True)
    apis = [("Context API", "/context — PAT-authenticated retrieval for apps"),
            ("MCP server", "/mcp — Model Context Protocol for AI agents"),
            ("Admin panel", "/admin — connector config, ingest, maintenance")]
    x = 0.9
    for name, sub in apis:
        rect(s, x, 4.65, 3.75, 1.35, fill=MIST, line_color=LINE,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + 0.3, 4.82, 3.2, 0.5, name, size=16, color=TEAL_DK, bold=True)
        txt(s, x + 0.3, 5.3, 3.25, 0.6, sub, size=12, color=SLATE,
            line_spacing=1.05)
        x += 3.95
    txt(s, 0.72, 6.35, 12, 0.4,
        "Webhooks + delta sync keep connectors current; per-participant ACLs "
        "carry through to query-time enforcement.", size=12, color=SLATE,
        italic=True)
    footer(s, 9, LABEL)

    # 10 — Deployment & observability
    s = blank(prs); bg(s)
    header(s, "Deployment & observability", kicker="Runs on Azure")
    # pipeline
    pipe = ["Local amd64 build", "Azure Container Registry", "Container Apps rollout",
            "Health verify"]
    x = 0.9
    for i, p in enumerate(pipe):
        c = rect(s, x, 2.35, 2.7, 1.0, fill=WHITE, line_color=AZURE, line_w=1.3,
                 shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        fill_text(c, p, size=13, color=NAVY, bold=True)
        if i < 3:
            arrow(s, x + 2.7, 2.85, 0.28, color=AZURE)
        x += 2.98
    # two columns: identity / observability
    rect(s, 0.9, 3.85, 5.7, 2.5, fill=MIST, line_color=LINE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 1.2, 4.0, 5.2, 0.4, "Identity & secrets", size=15, color=NAVY, bold=True)
    bullets(s, 1.2, 4.55, 5.2, [
        ("Managed identity", "DefaultAzureCredential, no secrets in code"),
        ("Azure Key Vault", "only what MI can't cover (OBO secret)"),
        ("Microsoft Entra", "SSO + JWT validation, PAT for apps"),
    ], size=12.5, gap=0.56)
    rect(s, 6.8, 3.85, 5.7, 2.5, fill=MIST, line_color=LINE,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, 7.1, 4.0, 5.2, 0.4, "Observability", size=15, color=NAVY, bold=True)
    bullets(s, 7.1, 4.55, 5.2, [
        ("OpenTelemetry", "one trace per request"),
        ("Application Insights", "signal breakdown as span attributes"),
        ("Log Analytics", "centralized Container Apps logs"),
    ], size=12.5, gap=0.56)
    txt(s, 0.72, 6.55, 12, 0.4,
        "Region: centralindia · Azure OpenAI in swedencentral. Bicep/APIM "
        "deferred to Phase 5.", size=12, color=SLATE, italic=True)
    footer(s, 10, LABEL)

    # 11 — What's next
    s = blank(prs); bg(s, NAVY)
    rect(s, 0, 0, 13.333, 0.16, fill=AZURE)
    txt(s, 0.9, 1.2, 11.5, 0.5, "PHASE 5 — INFRA HARDENING", size=14,
        color=AZURE_LT, bold=True)
    txt(s, 0.85, 1.8, 11.6, 1.0, "What's next", size=40, color=WHITE, bold=True,
        font=FONT_LT)
    nexts = [
        ("Per-user OBO", "delegated on-behalf-of tokens for Live Fetch"),
        ("APIM gateway", "auth, rate-limit, WAF in front of the API"),
        ("Event-driven ingest", "Event Hubs ingestion path"),
        ("Per-tenant isolation", "dedicated index per tenant"),
        ("JWKS caching", "token-validation performance"),
        ("Bicep IaC", "replace az-CLI provisioning"),
    ]
    x = 0.9
    y = 3.0
    for i, (name, sub) in enumerate(nexts):
        col = x + (i % 3) * 4.05
        row = y + (i // 3) * 1.5
        rect(s, col, row, 3.8, 1.3, fill=RGBColor(0x0E, 0x29, 0x4F),
             line_color=AZURE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, col + 0.28, row + 0.18, 3.3, 0.4, name, size=15, color=CYAN,
            bold=True)
        txt(s, col + 0.28, row + 0.62, 3.3, 0.6, sub, size=11.5, color=WHITE,
            line_spacing=1.05)
    txt(s, 0.9, 6.75, 11.5, 0.4,
        "Clean module boundaries today → mechanical lift to microservices "
        "tomorrow.", size=13, color=AZURE_LT, italic=True)

    out = os.path.join(os.path.dirname(__file__), "SubStrateOS_Engineering.pptx")
    prs.save(out)
    return out


if __name__ == "__main__":
    b = build_business()
    e = build_engineering()
    print("wrote:", b)
    print("wrote:", e)
